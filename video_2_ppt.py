"""
视频转PPT工具
从视频中提取PPT幻灯片并生成PPT文件
"""

import cv2
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io
import time


class VideoToPPT:
    """视频转PPT转换器"""
    
    def __init__(self, video_path, output_path="output.pptx", 
                 similarity_threshold=0.95, min_interval=2.0):
        """
        初始化转换器
        
        Args:
            video_path: 视频文件路径
            output_path: 输出PPT文件路径
            similarity_threshold: 相似度阈值（0-1），超过此值认为是同一张幻灯片
            min_interval: 最小间隔时间（秒），两张幻灯片之间的最小时间间隔
        """
        self.video_path = Path(video_path)
        self.output_path = Path(output_path)
        self.similarity_threshold = similarity_threshold
        self.min_interval = min_interval
        self.slides = []
        self.temp_dir = Path("temp_slides")
        self.temp_dir.mkdir(exist_ok=True)
    
    def extract_frames(self, interval=1.0):
        """
        从视频中提取帧
        
        Args:
            interval: 提取间隔（秒）
            
        Returns:
            list: 提取的帧列表 [(timestamp, frame), ...]
        """
        print(f"正在打开视频: {self.video_path}")
        
        cap = cv2.VideoCapture(str(self.video_path))
        
        if not cap.isOpened():
            raise ValueError(f"无法打开视频文件: {self.video_path}")
        
        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps
        
        print(f"视频信息:")
        print(f"  - 帧率: {fps:.2f} FPS")
        print(f"  - 总帧数: {total_frames}")
        print(f"  - 时长: {duration:.2f} 秒")
        
        frames = []
        frame_interval = int(fps * interval)
        
        print(f"\n开始提取帧，间隔: {interval} 秒...")
        
        frame_count = 0
        while True:
            ret, frame = cap.read()
            
            if not ret:
                break
            
            if frame_count % frame_interval == 0:
                timestamp = frame_count / fps
                frames.append((timestamp, frame))
                print(f"\r已提取 {len(frames)} 帧", end='', flush=True)
            
            frame_count += 1
        
        cap.release()
        print(f"\n提取完成，共 {len(frames)} 帧")
        
        return frames
    
    def calculate_similarity(self, frame1, frame2):
        """
        计算两帧之间的相似度
        
        Args:
            frame1: 第一帧
            frame2: 第二帧
            
        Returns:
            float: 相似度（0-1）
        """
        # 转换为灰度图
        gray1 = cv2.cvtColor(frame1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(frame2, cv2.COLOR_BGR2GRAY)
        
        # 调整大小以加快计算速度
        size = (100, 100)
        gray1 = cv2.resize(gray1, size)
        gray2 = cv2.resize(gray2, size)
        
        # 计算结构相似性
        mse = np.mean((gray1.astype(float) - gray2.astype(float)) ** 2)
        
        if mse == 0:
            return 1.0
        
        # 将MSE转换为相似度
        similarity = 1.0 / (1.0 + mse)
        
        return similarity
    
    def detect_slides(self, frames):
        """
        检测幻灯片变化
        
        Args:
            frames: 帧列表 [(timestamp, frame), ...]
            
        Returns:
            list: 幻灯片列表 [(timestamp, frame), ...]
        """
        print("\n开始检测幻灯片变化...")
        
        if not frames:
            return []
        
        slides = []
        last_frame = None
        last_timestamp = 0
        
        for i, (timestamp, frame) in enumerate(frames):
            print(f"\r处理进度: {i+1}/{len(frames)}", end='', flush=True)
            
            if last_frame is None:
                slides.append((timestamp, frame))
                last_frame = frame
                last_timestamp = timestamp
                continue
            
            # 检查时间间隔
            time_diff = timestamp - last_timestamp
            if time_diff < self.min_interval:
                continue
            
            # 计算相似度
            similarity = self.calculate_similarity(last_frame, frame)
            
            # 如果相似度低于阈值，认为是新的幻灯片
            if similarity < self.similarity_threshold:
                slides.append((timestamp, frame))
                last_frame = frame
                last_timestamp = timestamp
        
        print(f"\n检测完成，共找到 {len(slides)} 张幻灯片")
        
        return slides
    
    def save_slides_as_images(self, slides):
        """
        将幻灯片保存为图片
        
        Args:
            slides: 幻灯片列表
            
        Returns:
            list: 图片路径列表
        """
        print("\n保存幻灯片为图片...")
        
        image_paths = []
        
        for i, (timestamp, frame) in enumerate(slides):
            image_path = self.temp_dir / f"slide_{i+1:03d}.png"
            cv2.imwrite(str(image_path), frame)
            image_paths.append(image_path)
            print(f"\r已保存 {i+1}/{len(slides)} 张幻灯片", end='', flush=True)
        
        print(f"\n图片保存完成")
        
        return image_paths
    
    def create_ppt(self, image_paths):
        """
        创建PPT文件
        
        Args:
            image_paths: 图片路径列表
        """
        print(f"\n创建PPT文件: {self.output_path}")
        
        prs = Presentation()
        
        # 设置幻灯片大小为16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        for i, image_path in enumerate(image_paths):
            # 创建空白幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加图片
            img = Image.open(image_path)
            img_width, img_height = img.size
            
            # 计算图片在幻灯片中的位置和大小
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # 保持宽高比
            if img_width / img_height > slide_width / slide_height:
                width = slide_width
                height = width * img_height / img_width
            else:
                height = slide_height
                width = height * img_width / img_height
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            slide.shapes.add_picture(str(image_path), left, top, width, height)
            
            # 添加页码
            textbox = slide.shapes.add_textbox(Inches(9), Inches(5.2), Inches(0.5), Inches(0.3))
            text_frame = textbox.text_frame
            text_frame.text = str(i + 1)
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            
            print(f"\r已创建 {i+1}/{len(image_paths)} 张幻灯片", end='', flush=True)
        
        prs.save(str(self.output_path))
        print(f"\nPPT文件创建完成")
    
    def convert(self, extract_interval=1.0):
        """
        执行完整的转换流程
        
        Args:
            extract_interval: 提取帧的间隔（秒）
            
        Returns:
            bool: 是否转换成功
        """
        try:
            start_time = time.time()
            
            # 提取帧
            frames = self.extract_frames(extract_interval)
            
            if not frames:
                print("未提取到任何帧")
                return False
            
            # 检测幻灯片
            slides = self.detect_slides(frames)
            
            if not slides:
                print("未检测到任何幻灯片")
                return False
            
            # 保存为图片
            image_paths = self.save_slides_as_images(slides)
            
            # 创建PPT
            self.create_ppt(image_paths)
            
            end_time = time.time()
            elapsed_time = end_time - start_time
            
            print(f"\n转换完成！")
            print(f"耗时: {elapsed_time:.2f} 秒")
            print(f"输出文件: {self.output_path}")
            
            return True
            
        except Exception as e:
            print(f"\n转换过程出错: {str(e)}")
            import traceback
            traceback.print_exc()
            return False
    
    def cleanup(self):
        """清理临时文件"""
        try:
            for file in self.temp_dir.glob("*.png"):
                file.unlink()
            self.temp_dir.rmdir()
            print("临时文件已清理")
        except Exception as e:
            print(f"清理临时文件失败: {str(e)}")


def main():
    """主函数"""
    print("=" * 60)
    print("视频转PPT工具")
    print("=" * 60)
    
    video_path = "downloads\\RAG架构与实践指南v2.mp4"
    
    if not video_path:
        print("错误: 必须提供视频文件路径")
        return
    
    if not Path(video_path).exists():
        print(f"错误: 视频文件不存在: {video_path}")
        return
    
    output_path = input("请输入输出PPT文件路径 (直接回车使用默认): ").strip()
    if not output_path:
        output_path = "RAG架构与实践指南v2.pptx"
    
    similarity_threshold = input("请输入相似度阈值 (0-1, 默认 0.95): ").strip()
    if similarity_threshold:
        try:
            similarity_threshold = float(similarity_threshold)
            similarity_threshold = max(0.0, min(1.0, similarity_threshold))
        except ValueError:
            print("相似度阈值格式错误，使用默认值 0.95")
            similarity_threshold = 0.95
    else:
        similarity_threshold = 0.95
    
    min_interval = input("请输入最小间隔时间（秒，默认 2.0）: ").strip()
    if min_interval:
        try:
            min_interval = float(min_interval)
        except ValueError:
            print("最小间隔时间格式错误，使用默认值 2.0")
            min_interval = 2.0
    else:
        min_interval = 2.0
    
    extract_interval = input("请输入帧提取间隔（秒，默认 60.0）: ").strip()
    if extract_interval:
        try:
            extract_interval = float(extract_interval)
        except ValueError:
            print("帧提取间隔格式错误，使用默认值 60.0")
            extract_interval = 60.0
    else:
        extract_interval = 60.0
    
    cleanup = input("转换完成后是否清理临时文件? (y/n, 默认 y): ").strip().lower()
    cleanup = cleanup != 'n'
    
    print("\n开始转换...")
    print("-" * 60)
    
    converter = VideoToPPT(
        video_path=video_path,
        output_path=output_path,
        similarity_threshold=similarity_threshold,
        min_interval=min_interval
    )
    
    success = converter.convert(extract_interval)
    
    print("-" * 60)
    
    if success:
        print(f"✓ 转换成功！PPT文件保存在: {Path(output_path).absolute()}")
        
        if cleanup:
            converter.cleanup()
    else:
        print("✗ 转换失败，请检查视频文件和参数设置")


if __name__ == "__main__":
    main()
