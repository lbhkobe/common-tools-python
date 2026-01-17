"""
图片转PPT工具
将文件夹下的所有图片转换为PPT文件
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import os


class ImageToPPT:
    """图片转PPT转换器"""
    
    def __init__(self, image_folder, output_path="output.pptx", 
                 slide_width=10, slide_height=5.625):
        """
        初始化转换器
        
        Args:
            image_folder: 图片文件夹路径
            output_path: 输出PPT文件路径
            slide_width: 幻灯片宽度（英寸）
            slide_height: 幻灯片高度（英寸）
        """
        self.image_folder = Path(image_folder)
        self.output_path = Path(output_path)
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.supported_formats = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
    
    def get_image_files(self):
        """
        获取文件夹中的所有图片文件
        
        Returns:
            list: 图片文件路径列表（按文件名排序）
        """
        if not self.image_folder.exists():
            raise ValueError(f"文件夹不存在: {self.image_folder}")
        
        image_files = []
        
        for file in self.image_folder.iterdir():
            if file.is_file() and file.suffix.lower() in self.supported_formats:
                image_files.append(file)
        
        # 按文件名排序
        image_files.sort(key=lambda x: x.name)
        
        return image_files
    
    def create_ppt(self, image_files, add_page_numbers=True):
        """
        创建PPT文件
        
        Args:
            image_files: 图片文件路径列表
            add_page_numbers: 是否添加页码
        """
        print(f"\n创建PPT文件: {self.output_path}")
        
        prs = Presentation()
        
        # 设置幻灯片大小
        prs.slide_width = Inches(self.slide_width)
        prs.slide_height = Inches(self.slide_height)
        
        for i, image_path in enumerate(image_files):
            # 创建空白幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加图片
            img = Image.open(image_path)
            img_width, img_height = img.size
            
            # 计算图片在幻灯片中的位置和大小
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # 保持宽高比
            if img_width / img_height > slide_width / slide_height:
                width = slide_width
                height = width * img_height / img_width
            else:
                height = slide_height
                width = height * img_width / img_height
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            slide.shapes.add_picture(str(image_path), left, top, width, height)
            
            # 添加页码
            if add_page_numbers:
                textbox = slide.shapes.add_textbox(
                    Inches(self.slide_width - 0.8), 
                    Inches(self.slide_height - 0.4), 
                    Inches(0.5), 
                    Inches(0.3)
                )
                text_frame = textbox.text_frame
                text_frame.text = str(i + 1)
                text_frame.paragraphs[0].font.size = Pt(12)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            
            print(f"\r已添加 {i+1}/{len(image_files)} 张图片", end='', flush=True)
        
        prs.save(str(self.output_path))
        print(f"\nPPT文件创建完成")
    
    def convert(self, add_page_numbers=True):
        """
        执行完整的转换流程
        
        Args:
            add_page_numbers: 是否添加页码
            
        Returns:
            bool: 是否转换成功
        """
        try:
            # 获取图片文件
            image_files = self.get_image_files()
            
            if not image_files:
                print(f"文件夹中没有找到支持的图片文件")
                print(f"支持的格式: {', '.join(self.supported_formats)}")
                return False
            
            print(f"找到 {len(image_files)} 张图片")
            print(f"图片列表:")
            for i, img_file in enumerate(image_files[:10]):
                print(f"  {i+1}. {img_file.name}")
            if len(image_files) > 10:
                print(f"  ... 还有 {len(image_files) - 10} 张图片")
            
            # 创建PPT
            self.create_ppt(image_files, add_page_numbers)
            
            print(f"\n转换完成！")
            print(f"输出文件: {self.output_path}")
            print(f"幻灯片数量: {len(image_files)}")
            
            return True
            
        except Exception as e:
            print(f"\n转换过程出错: {str(e)}")
            import traceback
            traceback.print_exc()
            return False


def main():
    """主函数"""
    print("=" * 60)
    print("图片转PPT工具")
    print("=" * 60)
    
    image_folder = input("请输入图片文件夹路径 (直接回车使用当前目录): ").strip()
    if not image_folder:
        image_folder = "."
    
    if not Path(image_folder).exists():
        print(f"错误: 文件夹不存在: {image_folder}")
        return
    
    output_path = input("请输入输出PPT文件路径 (直接回车使用默认): ").strip()
    if not output_path:
        output_path = "output.pptx"
    
    print("\n请选择幻灯片尺寸:")
    print("1. 16:9 (推荐)")
    print("2. 4:3")
    
    size_choice = input("\n请选择 (1/2, 默认 1): ").strip()
    
    if size_choice == "2":
        slide_width = 10
        slide_height = 7.5
    else:
        slide_width = 10
        slide_height = 5.625
    
    add_page_numbers = input("是否添加页码? (y/n, 默认 y): ").strip().lower()
    add_page_numbers = add_page_numbers != 'n'
    
    print("\n开始转换...")
    print("-" * 60)
    
    converter = ImageToPPT(
        image_folder=image_folder,
        output_path=output_path,
        slide_width=slide_width,
        slide_height=slide_height
    )
    
    success = converter.convert(add_page_numbers)
    
    print("-" * 60)
    
    if success:
        print(f"✓ 转换成功！PPT文件保存在: {Path(output_path).absolute()}")
    else:
        print("✗ 转换失败，请检查文件夹和图片文件")


if __name__ == "__main__":
    main()
